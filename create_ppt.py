from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(1))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.alignment = PP_ALIGN.CENTER

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)
    
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    
    content_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(12), Inches(5.5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = content_frame.paragraphs[0]
        else:
            p = content_frame.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(22)
        p.space_before = Pt(12)

add_title_slide(prs, "使用 OpenCode 开发数据可视化模块", "技术分享汇报")

add_content_slide(prs, "项目概述", [
    "模块功能：高校社交媒体数据分析和可视化",
    "技术栈：Python + Pandas + Matplotlib",
    "数据内容：10所高校的校园文化、学术科研类文章",
    "可视化效果：3D柱状图、饼图等多维度展示"
])

add_content_slide(prs, "OpenCode 使用技巧", [
    "代码搜索：快速定位相关函数和类的实现",
    "文件读取：直接查看源码，无需跳转IDE",
    "多文件编辑：并行读取和修改多个文件",
    "任务拆分：复杂任务交给子agent处理",
    "Web搜索：遇到问题时快速查阅文档"
])

add_content_slide(prs, "问题与解决方案", [
    "3D图表坐标轴标签重叠 → 调整字体大小和旋转角度",
    "中文字体显示问题 → 配置 SimHei 等中文字体",
    "代码逻辑复杂 → 使用子agent分段处理",
    "初次使用不熟悉 → 从简单任务开始练习"
])

add_content_slide(prs, "使用心得", [
    "效率提升：减少上下文切换，快速完成开发",
    "代码质量：规范的代码风格和最佳实践",
    "学习曲线：上手简单，但需要熟悉prompt技巧",
    "适用场景：适合快速原型开发和小模块实现",
    "局限性：复杂业务逻辑仍需人工把控"
])

add_content_slide(prs, "总结与建议", [
    "OpenCode 是强大的AI辅助编程工具",
    "善于利用搜索和子任务功能提高效率",
    "注意验证AI生成的代码逻辑正确性",
    "建议团队推广，提升整体开发效率",
    "持续关注工具更新，获取新功能"
])

add_title_slide(prs, "谢谢！", "欢迎交流讨论")

prs.save('opencode技术分享.pptx')
print("PPT已生成: opencode技术分享.pptx")
